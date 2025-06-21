from typing import List, BinaryIO
from pptx import Presentation
from io import BytesIO
from .base import BaseParser

class PptxParser(BaseParser):
    """Parser for Microsoft PowerPoint presentations."""
    
    def parse(self, file: BinaryIO) -> List[str]:
        """
        Parse a PowerPoint presentation and return a list of text chunks.
        
        Args:
            file (BinaryIO): PowerPoint file object
            
        Returns:
            List[str]: List of text chunks
        """
        # Create presentation object
        prs = Presentation(BytesIO(file.read()))
        
        # Extract text from slides
        text_content = []
        
        for slide in prs.slides:
            slide_text = []
            
            # Get text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                
                # Get text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_text.append(' | '.join(row_text))
            
            if slide_text:
                # Join all text from the slide
                text_content.append('\n'.join(slide_text))
        
        # Join all slides and chunk the text
        full_text = '\n\n'.join(text_content)
        return self.chunk_text(self.clean_text(full_text)) 